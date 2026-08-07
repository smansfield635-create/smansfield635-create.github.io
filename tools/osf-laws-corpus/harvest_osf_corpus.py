#!/usr/bin/env python3
"""Harvest a public OSF user corpus into deterministic, reviewable GitHub artifacts.

This script does not validate scientific claims. It inventories public OSF nodes,
registrations, files, and extractable text, then emits provenance-preserving
records for later Laws Chamber review.
"""

from __future__ import annotations

import argparse
import csv
import hashlib
import io
import json
import os
import re
import shutil
import sys
import time
from dataclasses import dataclass, asdict
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Iterable
from urllib.parse import urlparse

import requests
from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

API_ROOT = "https://api.osf.io/v2"
SESSION = requests.Session()
SESSION.headers.update(
    {
        "Accept": "application/vnd.api+json",
        "User-Agent": "laws-osf-corpus-harvester/1.0 (+https://github.com/smansfield635-create/smansfield635-create.github.io)",
    }
)

TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".rst", ".csv", ".tsv", ".json", ".jsonl",
    ".yaml", ".yml", ".xml", ".html", ".htm", ".css", ".js", ".mjs", ".cjs",
    ".ts", ".tsx", ".jsx", ".py", ".r", ".sql", ".tex", ".bib", ".ini", ".cfg",
    ".toml", ".log",
}
EXTRACTABLE_EXTENSIONS = TEXT_EXTENSIONS | {".pdf", ".docx", ".xlsx", ".xlsm", ".pptx"}

LAW_TERMS = {
    "law": 5,
    "coherence": 5,
    "constraint": 4,
    "invariant": 4,
    "falsification": 5,
    "scientific method": 5,
    "quotient": 3,
    "propagation": 3,
    "failure": 2,
    "trajectory": 3,
    "transitory": 4,
    "cyclical": 3,
    "field phenomenon": 4,
    "instantiation": 4,
    "diagnostic": 3,
    "alignment": 2,
    "boundary": 3,
    "formalism": 3,
    "seven invariants": 5,
    "m256": 3,
    "pfq": 4,
    "ccis": 4,
}
EVIDENCE_TERMS = {
    "dataset": 3,
    "data": 1,
    "method": 2,
    "methods": 2,
    "result": 2,
    "results": 2,
    "sample": 2,
    "experiment": 3,
    "empirical": 4,
    "falsif": 4,
    "appendix": 2,
    "table": 1,
    "figure": 1,
    "code": 2,
    "validation": 3,
    "pilot": 3,
}

MAX_FILE_BYTES = 25 * 1024 * 1024
MAX_TEXT_CHARS_PER_FILE = 600_000
MAX_TEXT_CHARS_PER_NODE = 3_000_000
MAX_TOTAL_DOWNLOAD_BYTES = 750 * 1024 * 1024
REQUEST_TIMEOUT = 60


class HarvestError(RuntimeError):
    pass


@dataclass
class FileRecord:
    node_id: str
    provider: str
    name: str
    kind: str
    materialized_path: str | None
    size: int | None
    content_type: str | None
    date_modified: str | None
    download_url: str | None
    osf_file_id: str | None
    extraction_status: str
    extracted_chars: int
    sha256: str | None
    error: str | None


def utc_now() -> str:
    return datetime.now(timezone.utc).isoformat()


def safe_name(value: str, fallback: str = "untitled") -> str:
    value = re.sub(r"[^A-Za-z0-9._-]+", "-", value.strip()).strip("-.")
    return value[:120] or fallback


def request_json(url: str, *, retries: int = 4) -> dict[str, Any]:
    delay = 1.0
    last_error: Exception | None = None
    for attempt in range(retries):
        try:
            response = SESSION.get(url, timeout=REQUEST_TIMEOUT)
            if response.status_code in {429, 500, 502, 503, 504}:
                raise HarvestError(f"transient HTTP {response.status_code} for {url}")
            response.raise_for_status()
            return response.json()
        except Exception as exc:  # noqa: BLE001
            last_error = exc
            if attempt + 1 >= retries:
                break
            time.sleep(delay)
            delay *= 2
    raise HarvestError(f"failed JSON request {url}: {last_error}")


def paginate(url: str) -> list[dict[str, Any]]:
    out: list[dict[str, Any]] = []
    seen: set[str] = set()
    while url:
        if url in seen:
            raise HarvestError(f"pagination loop detected: {url}")
        seen.add(url)
        payload = request_json(url)
        data = payload.get("data", [])
        if isinstance(data, dict):
            out.append(data)
        else:
            out.extend(data)
        url = (payload.get("links") or {}).get("next") or ""
    return out


def relationship_href(resource: dict[str, Any], name: str) -> str | None:
    rel = ((resource.get("relationships") or {}).get(name) or {})
    links = rel.get("links") or {}
    related = links.get("related")
    if isinstance(related, dict):
        return related.get("href")
    if isinstance(related, str):
        return related
    return None


def relationship_id(resource: dict[str, Any], name: str) -> str | None:
    rel = ((resource.get("relationships") or {}).get(name) or {})
    data = rel.get("data")
    if isinstance(data, dict):
        return data.get("id")
    href = relationship_href(resource, name)
    if href:
        parts = [part for part in urlparse(href).path.split("/") if part]
        if parts:
            return parts[-1]
    return None


def clean_html(value: str | None) -> str:
    if not value:
        return ""
    return BeautifulSoup(value, "html.parser").get_text("\n", strip=True)


def normalize_node(resource: dict[str, Any], *, source: str) -> dict[str, Any]:
    attrs = resource.get("attributes") or {}
    node_id = resource.get("id")
    links = resource.get("links") or {}
    description = clean_html(attrs.get("description"))
    parent_id = relationship_id(resource, "parent")
    return {
        "id": node_id,
        "type": resource.get("type"),
        "title": attrs.get("title") or "",
        "description": description,
        "category": attrs.get("category"),
        "public": attrs.get("public"),
        "registration": attrs.get("registration"),
        "preprint": attrs.get("preprint"),
        "date_created": attrs.get("date_created"),
        "date_modified": attrs.get("date_modified"),
        "current_user_permissions": attrs.get("current_user_permissions"),
        "tags": attrs.get("tags") or [],
        "subjects": attrs.get("subjects") or [],
        "node_license": attrs.get("node_license"),
        "parent_id": parent_id,
        "is_component": bool(parent_id),
        "osf_url": links.get("html") or f"https://osf.io/{node_id}/",
        "api_url": links.get("self") or f"{API_ROOT}/nodes/{node_id}/",
        "source_collection": source,
    }


def enumerate_user_resources(user_id: str) -> tuple[dict[str, dict[str, Any]], list[str]]:
    resources: dict[str, dict[str, Any]] = {}
    errors: list[str] = []
    endpoints = [
        ("user_nodes", f"{API_ROOT}/users/{user_id}/nodes/?page%5Bsize%5D=100"),
        ("user_registrations", f"{API_ROOT}/users/{user_id}/registrations/?page%5Bsize%5D=100"),
    ]
    for label, endpoint in endpoints:
        try:
            for resource in paginate(endpoint):
                rid = resource.get("id")
                if rid:
                    resources[rid] = resource | {"_source_collection": label}
        except Exception as exc:  # noqa: BLE001
            errors.append(f"{label}: {exc}")

    queue = list(resources.values())
    visited_children: set[str] = set()
    while queue:
        resource = queue.pop(0)
        rid = resource.get("id")
        if not rid or rid in visited_children:
            continue
        visited_children.add(rid)
        href = relationship_href(resource, "children")
        if not href:
            continue
        try:
            for child in paginate(href):
                cid = child.get("id")
                if cid and cid not in resources:
                    child["_source_collection"] = "recursive_child"
                    resources[cid] = child
                    queue.append(child)
        except Exception as exc:  # noqa: BLE001
            errors.append(f"children:{rid}: {exc}")
    return resources, errors


def list_node_files(node_id: str) -> tuple[list[dict[str, Any]], list[str]]:
    files: list[dict[str, Any]] = []
    errors: list[str] = []
    try:
        providers = paginate(f"{API_ROOT}/nodes/{node_id}/files/?page%5Bsize%5D=100")
    except Exception as exc:  # noqa: BLE001
        return [], [f"providers:{node_id}: {exc}"]

    for provider in providers:
        provider_name = (provider.get("attributes") or {}).get("name") or provider.get("id") or "unknown"
        href = relationship_href(provider, "files")
        if not href:
            continue
        queue = [href]
        seen: set[str] = set()
        while queue:
            current = queue.pop(0)
            if current in seen:
                continue
            seen.add(current)
            try:
                entries = paginate(current)
            except Exception as exc:  # noqa: BLE001
                errors.append(f"files:{node_id}:{provider_name}: {exc}")
                continue
            for entry in entries:
                entry["_provider"] = provider_name
                files.append(entry)
                attrs = entry.get("attributes") or {}
                if attrs.get("kind") == "folder":
                    child_href = relationship_href(entry, "files")
                    if child_href:
                        queue.append(child_href)
    return files, errors


def read_text_bytes(data: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "utf-16", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def extract_text(name: str, data: bytes) -> str:
    ext = Path(name).suffix.lower()
    if ext in TEXT_EXTENSIONS:
        text = read_text_bytes(data)
        if ext in {".html", ".htm", ".xml"}:
            text = BeautifulSoup(text, "html.parser").get_text("\n", strip=True)
        return text
    if ext == ".pdf":
        reader = PdfReader(io.BytesIO(data))
        pages: list[str] = []
        for idx, page in enumerate(reader.pages, start=1):
            pages.append(f"\n\n--- PAGE {idx} ---\n\n{page.extract_text() or ''}")
        return "".join(pages)
    if ext == ".docx":
        doc = Document(io.BytesIO(data))
        chunks = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                chunks.append("\t".join(cell.text for cell in row.cells))
        return "\n".join(chunks)
    if ext in {".xlsx", ".xlsm"}:
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        chunks: list[str] = []
        for ws in wb.worksheets:
            chunks.append(f"\n--- SHEET: {ws.title} ---")
            for row in ws.iter_rows(values_only=True):
                chunks.append("\t".join("" if v is None else str(v) for v in row))
        return "\n".join(chunks)
    if ext == ".pptx":
        prs = Presentation(io.BytesIO(data))
        chunks: list[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            chunks.append(f"\n--- SLIDE {idx} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    chunks.append(shape.text)
        return "\n".join(chunks)
    raise ValueError(f"unsupported extension: {ext or '(none)'}")


def score_terms(text: str, terms: dict[str, int]) -> tuple[int, list[str]]:
    lower = text.lower()
    hits: list[str] = []
    score = 0
    for term, weight in terms.items():
        if term in lower:
            hits.append(term)
            score += weight
    return score, hits


def classify_relevance(title: str, description: str, extracted: str, files: list[FileRecord]) -> dict[str, Any]:
    combined = "\n".join([title, description, extracted[:1_500_000]])
    law_score, law_hits = score_terms(combined, LAW_TERMS)
    evidence_score, evidence_hits = score_terms(combined, EVIDENCE_TERMS)
    extractable_count = sum(1 for f in files if f.extraction_status in {"EXTRACTED", "EXTRACTED_TRUNCATED"})
    has_data_file = any(Path(f.name).suffix.lower() in {".csv", ".tsv", ".xlsx", ".xlsm", ".json", ".jsonl"} for f in files)
    has_code_file = any(Path(f.name).suffix.lower() in {".py", ".r", ".js", ".mjs", ".ts", ".sql"} for f in files)
    if law_score >= 18:
        band = "CORE_LAWS_CANDIDATE"
    elif law_score >= 9:
        band = "STRONG_SUPPORTING_CANDIDATE"
    elif law_score >= 4:
        band = "APPLIED_OR_CONTEXTUAL_CANDIDATE"
    else:
        band = "LOW_DIRECT_LAWS_SIGNAL"
    return {
        "candidate_band": band,
        "law_relevance_score": law_score,
        "law_term_hits": law_hits,
        "evidence_signal_score": evidence_score,
        "evidence_term_hits": evidence_hits,
        "extractable_file_count": extractable_count,
        "has_data_file": has_data_file,
        "has_code_file": has_code_file,
        "interpretation_boundary": "Heuristic discovery signal only; not a validation, empirical-support, or admission decision.",
    }


def harvest_node(node: dict[str, Any], output_root: Path, budget: dict[str, int]) -> dict[str, Any]:
    node_id = node["id"]
    node_dir = output_root / "projects" / node_id
    node_dir.mkdir(parents=True, exist_ok=True)
    raw_entries, file_errors = list_node_files(node_id)
    file_records: list[FileRecord] = []
    extracted_sections: list[str] = []
    project_chars = 0

    for entry in raw_entries:
        attrs = entry.get("attributes") or {}
        links = entry.get("links") or {}
        name = attrs.get("name") or entry.get("id") or "unnamed"
        kind = attrs.get("kind") or "unknown"
        size = attrs.get("size")
        download_url = links.get("download")
        record = FileRecord(
            node_id=node_id,
            provider=entry.get("_provider") or "unknown",
            name=name,
            kind=kind,
            materialized_path=attrs.get("materialized_path") or attrs.get("path"),
            size=size if isinstance(size, int) else None,
            content_type=attrs.get("content_type") or attrs.get("contentType"),
            date_modified=attrs.get("date_modified") or attrs.get("modified"),
            download_url=download_url,
            osf_file_id=entry.get("id"),
            extraction_status="NOT_ATTEMPTED",
            extracted_chars=0,
            sha256=None,
            error=None,
        )
        if kind != "file":
            record.extraction_status = "FOLDER"
            file_records.append(record)
            continue
        ext = Path(name).suffix.lower()
        if ext not in EXTRACTABLE_EXTENSIONS:
            record.extraction_status = "UNSUPPORTED_FORMAT"
            file_records.append(record)
            continue
        if not download_url:
            record.extraction_status = "NO_DOWNLOAD_URL"
            file_records.append(record)
            continue
        if record.size and record.size > MAX_FILE_BYTES:
            record.extraction_status = "SKIPPED_FILE_SIZE_LIMIT"
            file_records.append(record)
            continue
        if budget["downloaded"] >= MAX_TOTAL_DOWNLOAD_BYTES:
            record.extraction_status = "SKIPPED_TOTAL_SIZE_LIMIT"
            file_records.append(record)
            continue
        try:
            response = SESSION.get(download_url, timeout=REQUEST_TIMEOUT)
            response.raise_for_status()
            data = response.content
            budget["downloaded"] += len(data)
            record.sha256 = hashlib.sha256(data).hexdigest()
            if len(data) > MAX_FILE_BYTES:
                record.extraction_status = "SKIPPED_FILE_SIZE_LIMIT_AFTER_DOWNLOAD"
                file_records.append(record)
                continue
            text = extract_text(name, data).replace("\x00", "")
            remaining_project = MAX_TEXT_CHARS_PER_NODE - project_chars
            allowed = min(MAX_TEXT_CHARS_PER_FILE, max(0, remaining_project))
            if allowed <= 0:
                record.extraction_status = "SKIPPED_PROJECT_TEXT_LIMIT"
                file_records.append(record)
                continue
            truncated = len(text) > allowed
            text = text[:allowed]
            project_chars += len(text)
            record.extracted_chars = len(text)
            record.extraction_status = "EXTRACTED_TRUNCATED" if truncated else "EXTRACTED"
            extracted_sections.append(
                f"\n\n# FILE: {name}\n\n"
                f"- OSF materialized path: `{record.materialized_path or ''}`\n"
                f"- SHA-256: `{record.sha256}`\n"
                f"- Extraction status: `{record.extraction_status}`\n\n"
                f"```text\n{text}\n```\n"
            )
        except Exception as exc:  # noqa: BLE001
            record.extraction_status = "EXTRACTION_FAILED"
            record.error = str(exc)[:1000]
        file_records.append(record)

    extracted_text = "".join(extracted_sections)
    relevance = classify_relevance(node["title"], node["description"], extracted_text, file_records)
    metadata = node | {
        "harvested_at": utc_now(),
        "file_inventory_count": len(file_records),
        "file_errors": file_errors,
        "relevance_discovery": relevance,
        "claim_boundary": "Corpus inventory only. Presence in this repository does not establish correctness, novelty, validation, replication, or Laws Chamber admission.",
    }
    (node_dir / "metadata.json").write_text(json.dumps(metadata, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    (node_dir / "files.json").write_text(json.dumps([asdict(f) for f in file_records], indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    extracted_header = (
        f"# Extracted text — {node['title']}\n\n"
        f"OSF project: {node['osf_url']}\n\n"
        "This is machine-extracted source text for discovery. Layout, equations, figures, and tables may be degraded. "
        "Consult the OSF source file before relying on any passage.\n"
    )
    (node_dir / "extracted-text.md").write_text(extracted_header + extracted_text, encoding="utf-8")
    return metadata


def build_index(output_root: Path, nodes: list[dict[str, Any]], discovery_errors: list[str], user_id: str, budget: dict[str, int]) -> None:
    ordered = sorted(
        nodes,
        key=lambda n: (
            -int((n.get("relevance_discovery") or {}).get("law_relevance_score", 0)),
            (n.get("title") or "").lower(),
        ),
    )
    manifest = {
        "corpus_id": "OSF_PUBLIC_LAWS_DISCOVERY_CORPUS_v1",
        "osf_user_id": user_id,
        "osf_profile_url": f"https://osf.io/{user_id}/",
        "generated_at": utc_now(),
        "public_resource_count": len(ordered),
        "root_project_count": sum(1 for n in ordered if not n.get("is_component")),
        "component_count": sum(1 for n in ordered if n.get("is_component")),
        "registration_count": sum(1 for n in ordered if n.get("registration") or n.get("type") == "registrations"),
        "downloaded_bytes": budget["downloaded"],
        "discovery_errors": discovery_errors,
        "claim_boundary": "Discovery corpus only; no claim upgrade or admission authority.",
        "projects": ordered,
    }
    (output_root / "manifest.json").write_text(json.dumps(manifest, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    with (output_root / "corpus.jsonl").open("w", encoding="utf-8") as handle:
        for node in ordered:
            handle.write(json.dumps(node, ensure_ascii=False) + "\n")

    lines = [
        "# OSF Public Laws Discovery Corpus",
        "",
        f"Profile: https://osf.io/{user_id}/",
        "",
        f"Generated: `{manifest['generated_at']}`",
        "",
        f"Resources harvested: **{len(ordered)}**",
        "",
        "## Evidence boundary",
        "",
        "This directory is a provenance-preserving discovery corpus. It does not validate the source projects, "
        "upgrade empirical claims, establish novelty, or authorize Laws Chamber admission. Every substantive use "
        "must return to the cited OSF project and source file.",
        "",
        "## Candidate index",
        "",
        "| Score | Candidate band | Resource | OSF | Extractable files | Data | Code |",
        "|---:|---|---|---|---:|:---:|:---:|",
    ]
    for node in ordered:
        rel = node.get("relevance_discovery") or {}
        title = (node.get("title") or "Untitled").replace("|", "\\|")
        lines.append(
            f"| {rel.get('law_relevance_score', 0)} | {rel.get('candidate_band', '')} | "
            f"[{title}](projects/{node['id']}/metadata.json) | "
            f"[OSF]({node['osf_url']}) | {rel.get('extractable_file_count', 0)} | "
            f"{'Y' if rel.get('has_data_file') else 'N'} | {'Y' if rel.get('has_code_file') else 'N'} |"
        )
    if discovery_errors:
        lines.extend(["", "## Harvest warnings", ""])
        lines.extend(f"- `{error}`" for error in discovery_errors)
    (output_root / "README.md").write_text("\n".join(lines) + "\n", encoding="utf-8")


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--user", required=True, help="Public OSF user identifier")
    parser.add_argument("--output", required=True, type=Path)
    args = parser.parse_args()

    output_root: Path = args.output
    projects_dir = output_root / "projects"
    if projects_dir.exists():
        shutil.rmtree(projects_dir)
    output_root.mkdir(parents=True, exist_ok=True)

    resources, discovery_errors = enumerate_user_resources(args.user)
    if not resources:
        raise HarvestError(
            "No public OSF resources were returned. Confirm the user identifier and public visibility. "
            f"Errors: {discovery_errors}"
        )

    normalized: list[dict[str, Any]] = []
    for resource in resources.values():
        source = resource.pop("_source_collection", "unknown")
        normalized.append(normalize_node(resource, source=source))

    budget = {"downloaded": 0}
    harvested: list[dict[str, Any]] = []
    for index, node in enumerate(sorted(normalized, key=lambda n: (n["title"].lower(), n["id"])), start=1):
        print(f"[{index}/{len(normalized)}] harvesting {node['id']} {node['title']}", flush=True)
        try:
            harvested.append(harvest_node(node, output_root, budget))
        except Exception as exc:  # noqa: BLE001
            discovery_errors.append(f"node:{node['id']}: {exc}")
            fallback = node | {
                "harvested_at": utc_now(),
                "file_inventory_count": 0,
                "file_errors": [str(exc)],
                "relevance_discovery": classify_relevance(node["title"], node["description"], "", []),
                "claim_boundary": "Corpus inventory only; node file harvest failed.",
            }
            node_dir = output_root / "projects" / node["id"]
            node_dir.mkdir(parents=True, exist_ok=True)
            (node_dir / "metadata.json").write_text(json.dumps(fallback, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
            (node_dir / "files.json").write_text("[]\n", encoding="utf-8")
            (node_dir / "extracted-text.md").write_text(f"# Extracted text — {node['title']}\n\nHarvest failed: {exc}\n", encoding="utf-8")
            harvested.append(fallback)

    build_index(output_root, harvested, discovery_errors, args.user, budget)
    print(json.dumps({"resources": len(harvested), "downloaded_bytes": budget["downloaded"], "errors": len(discovery_errors)}))
    return 0


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except Exception as exc:  # noqa: BLE001
        print(f"HARVEST_FAILED: {exc}", file=sys.stderr)
        raise
